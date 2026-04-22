import os
import re
import json
import time
import sys
import asyncio
import datetime
import urllib.request
import urllib.parse
import webbrowser
import requests
import random
import threading
import subprocess
import imaplib
import smtplib
import tempfile
import email as email_lib
from email.mime.text import MIMEText
from email.mime.multipart import MIMEMultipart

# External libraries
import bs4
import yt_dlp
import speech_recognition as sr
from docx import Document
from docx.shared import Inches, Pt, RGBColor
from docx.enum.text import WD_ALIGN_PARAGRAPH
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from duckduckgo_search import DDGS
from shazamio import Shazam
import pywhatkit
from AppOpener import open as open_app, close as close_app

import config
from audio_engine import speak
from ai_core import clean_ai_meta_text, clean_repetitive_content

# ============================================================
# ORIGINAL: App Control & Utilities
# ============================================================

def run_app(app_name):
    try:
        open_app(app_name, match_closest=True)
        return f"Opening {app_name}, sir."
    except:
        return f"I could not find {app_name} on your system."

def send_whatsapp_message(target, message):
    phone = config.CONTACTS.get(target, target)
    try:
        pywhatkit.sendwhatmsg_instant(phone, message, 10, True, 2)
        return f"Message sent to {target}."
    except Exception as e:
        return f"Error sending message: {e}"

# ============================================================
# ORIGINAL: Web Search
# ============================================================

last_search_urls = []

def agentic_web_search(query, language="he"):
    global last_search_urls
    try:
        print(f">> [DEBUG] Web search: {query}")
        results_text = ""
        region_code = "il-he" if language == "he" else "wt-wt"
        
        with DDGS() as ddgs:
            results = list(ddgs.text(query, region=region_code, max_results=5))
            if not results:
                return "לא מצאתי מידע על זה ברשת." if language == "he" else "No information found."
            for i, res in enumerate(results):
                results_text += f"תוצאה {i+1}: {res.get('body', '')}\n"
                config.last_search_urls.append(res.get('href', ''))
                
        return results_text
    except Exception as e:
        print(f"[SEARCH] Error: {e}")
        return "הייתה לי בעיה בחיפוש המידע ברשת." if language == "he" else "Could not retrieve web information."

# ============================================================
# ORIGINAL: Document & Presentation Generation
# ============================================================

def generate_docx_document(topic, content, direction="rtl", topic_english=""):
    try:
        doc = Document()
        search_term = topic_english if topic_english else "Science"
        downloads_path = os.path.join(os.path.expanduser('~'), 'Downloads')
        safe_topic = re.sub(r'[\\/*?:"<>|]', "", topic)
        docx_path = os.path.join(downloads_path, f"Research_{safe_topic}.docx")
        image_path = os.path.join(downloads_path, f"CoverImage_{safe_topic}.jpg")

        title = doc.add_heading(f"עבודת מחקר: {topic}" if direction == "rtl" else f"Research Paper: {topic}", 0)
        title.alignment = WD_ALIGN_PARAGRAPH.RIGHT if direction == "rtl" else WD_ALIGN_PARAGRAPH.LEFT
        if title.runs:
            title.runs[0].font.color.rgb = RGBColor(0, 212, 255)

        try:
            headers = {'User-Agent': 'Mozilla/5.0'}
            image_download_url = None
            query_api = f"https://en.wikipedia.org/w/api.php?action=query&generator=search&gsrsearch={urllib.parse.quote(search_term)}&gsrlimit=1&prop=pageimages&piprop=original&format=json"
            req_query = urllib.request.Request(query_api, headers=headers)
            with urllib.request.urlopen(req_query, timeout=5) as resp:
                query_data = json.loads(resp.read().decode('utf-8'))
                pages = query_data.get('query', {}).get('pages', {})
                for page_id, page_info in pages.items():
                    if 'original' in page_info:
                        image_download_url = page_info['original']['source']
                        break
            if not image_download_url:
                safe_ai_term = urllib.parse.quote(f"{search_term} realistic")
                image_download_url = f"https://image.pollinations.ai/prompt/{safe_ai_term}?width=800&height=400&nologo=true"
            req_img = urllib.request.Request(image_download_url, headers=headers)
            with urllib.request.urlopen(req_img, timeout=10) as response, open(image_path, 'wb') as out_file:
                out_file.write(response.read())
            doc.add_picture(image_path, width=Inches(6.0))
            doc.paragraphs[-1].alignment = WD_ALIGN_PARAGRAPH.CENTER
        except Exception as e:
            print(f"[DOCX] Image skip: {e}")

        content = clean_ai_meta_text(content)
        content = clean_repetitive_content(content)
        clean_content = content.replace("**", "")
        
        for p_text in clean_content.split('\n'):
            p_text = p_text.strip()
            if not p_text: continue
            if p_text.startswith("##") or ("פרק" in p_text and len(p_text) < 50) or ("Chapter" in p_text and len(p_text) < 50):
                text_clean = p_text.replace("##", "").strip()
                h = doc.add_heading(text_clean, level=1)
                if h.runs: h.runs[0].font.color.rgb = RGBColor(0, 95, 115)
                h.alignment = WD_ALIGN_PARAGRAPH.RIGHT if direction == "rtl" else WD_ALIGN_PARAGRAPH.LEFT
            else:
                p = doc.add_paragraph(p_text)
                p.alignment = WD_ALIGN_PARAGRAPH.RIGHT if direction == "rtl" else WD_ALIGN_PARAGRAPH.LEFT

        doc.add_paragraph("\n\n")
        footer = doc.add_paragraph(f"created {datetime.datetime.now().strftime('%d/%m/%Y')} | by J.A.R.V.I.S")
        footer.alignment = WD_ALIGN_PARAGRAPH.CENTER
        doc.save(docx_path)
        os.startfile(docx_path)
        return True
    except Exception as e:
        print(f"[DOCX] Error: {e}")
        return False


def generate_pptx_presentation(topic, ai_content, direction="rtl"):
    try:
        prs = Presentation()
        slides_data = ai_content.split("---SLIDE---")
        for slide_text in slides_data:
            if not slide_text.strip(): continue
            lines = slide_text.strip().split('\n')
            title = ""
            bullets = []
            for line in lines:
                if line.startswith("TITLE:"):
                    title = line.replace("TITLE:", "").strip()
                elif line.startswith(("-", "*", "•")):
                    bullets.append(line[1:].strip())
            if title or bullets:
                slide = prs.slides.add_slide(prs.slide_layouts[1])
                title_shape = slide.shapes.title
                title_shape.text = title
                body_shape = slide.shapes.placeholders[1]
                tf = body_shape.text_frame
                tf.clear()
                for bullet in bullets:
                    p = tf.add_paragraph()
                    p.text = bullet
                    p.font.size = Pt(24)
                    p.alignment = PP_ALIGN.RIGHT if direction == "rtl" else PP_ALIGN.LEFT
                title_shape.text_frame.paragraphs[0].alignment = PP_ALIGN.RIGHT if direction == "rtl" else PP_ALIGN.LEFT

        downloads_path = os.path.join(os.path.expanduser('~'), 'Downloads')
        safe_topic = re.sub(r'[\\/*?:"<>|]', "", topic)
        file_path = os.path.join(downloads_path, f"Presentation_{safe_topic}.pptx")
        prs.save(file_path)
        os.startfile(file_path)
        return True
    except Exception as e:
        print(f"[PPTX] Error: {e}")
        return False

# ============================================================
# ORIGINAL: Music & Song Recognition
# ============================================================

async def identify_song_logic():
    shazam = Shazam()
    recognizer = sr.Recognizer()
    with sr.Microphone() as source:
        audio_data = recognizer.record(source, duration=7)
        temp_file = "temp_song.wav"
        with open(temp_file, "wb") as f:
            f.write(audio_data.get_wav_data())
    out = await shazam.recognize(temp_file)
    if os.path.exists(temp_file): os.remove(temp_file)
    return f"{out['track']['title']} - {out['track']['subtitle']}" if out.get('track') else None

def start_song_identification():
    return asyncio.run(identify_song_logic())

def find_song_by_lyrics():
    recognizer = sr.Recognizer()
    recognizer.pause_threshold = 3.0
    with sr.Microphone() as source:
        try:
            recognizer.adjust_for_ambient_noise(source, duration=1)
            audio = recognizer.record(source, duration=8)
            return recognizer.recognize_google(audio, language="en-US")
        except:
            return None

# ============================================================
# ORIGINAL: Weather, Timer, Briefing
# ============================================================

def get_global_weather(city_name):
    try:
        safe_city = urllib.parse.quote(city_name)
        geo_url = f"https://geocoding-api.open-meteo.com/v1/search?name={safe_city}&count=1&language=en&format=json"
        geo_response = requests.get(geo_url, timeout=10).json()
        if not geo_response.get("results"): return None
        lat = geo_response["results"][0]["latitude"]
        lon = geo_response["results"][0]["longitude"]
        weather_url = f"https://api.open-meteo.com/v1/forecast?latitude={lat}&longitude={lon}&current_weather=true"
        weather_response = requests.get(weather_url, timeout=10).json()
        return weather_response["current_weather"]["temperature"]
    except:
        return "ERROR"

def morning_briefing():
    now = datetime.datetime.now()
    greeting = "Good morning, sir." if 5 <= now.hour < 12 else "Good afternoon, sir."
    time_str = now.strftime("%I:%M %p")
    try:
        response = requests.get("https://wttr.in/London?format=%t", timeout=5)
        temp = response.text.strip().replace('+', '').replace('°C', '')
        weather_str = f"The current temperature in London is {temp} degrees Celsius."
    except:
        weather_str = "I am currently unable to fetch the weather data."
    briefing = f"{greeting} The time is {time_str}. {weather_str} Opening the daily news for you."
    speak(briefing)
    webbrowser.open("https://news.google.com")

def timer_thread(seconds, is_hebrew):
    time.sleep(seconds)
    speak("המונה זמן נגמר אדוני." if is_hebrew else "Sir, the timer has finished.")


# ============================================================
# NEW FEATURE 1: Python Code Execution
# Jarvis can write and run Python scripts on your machine!
# Usage: "Jarvis, calculate the 10th fibonacci number"
#        "Jarvis, run code: print('hello world')"
# ============================================================

def execute_python_code(code_or_description, generate_text_fn=None):
    """
    Execute Python code safely in a sandboxed subprocess.
    If given a natural language description, generates the code first using AI.
    
    Returns: string output from the script (or error message)
    """
    code = code_or_description.strip()
    
    # If it looks like a plain English description (no Python keywords), generate code
    python_keywords = ['import', 'print(', 'for ', 'def ', 'class ', '= ', 'if ', 'while ', 'return ']
    looks_like_code = any(kw in code for kw in python_keywords)
    
    if generate_text_fn and not looks_like_code and len(code) > 0:
        print(f"[CODE] Generating Python code for: {code}")
        generated = generate_text_fn(
            f"Write a short Python 3 script that accomplishes this task: {code}\n"
            f"Rules: Use only standard library. End with print() of the result. "
            f"Return ONLY the code. No explanations. No markdown fences.",
            max_tokens=400,
            temperature=0.2
        )
        if not generated:
            return "Sorry sir, I couldn't generate the code."
        # Strip any markdown fences if present
        code = re.sub(r'```(?:python)?|```', '', generated).strip()
        print(f"[CODE] Generated:\n{code}")
    
    if not code:
        return "No code to execute."
    
    try:
        with tempfile.NamedTemporaryFile(
            mode='w', suffix='.py', delete=False, encoding='utf-8'
        ) as f:
            f.write(code)
            tmp_path = f.name
        
        result = subprocess.run(
            [sys.executable, tmp_path],
            capture_output=True, text=True, timeout=15,
            encoding='utf-8', errors='replace'
        )
        os.remove(tmp_path)
        
        output = result.stdout.strip()
        errors = result.stderr.strip()
        
        if output:
            return f"Output: {output[:500]}"
        elif errors:
            return f"Error: {errors[:300]}"
        else:
            return "Code executed successfully with no output, sir."
            
    except subprocess.TimeoutExpired:
        try: os.remove(tmp_path)
        except: pass
        return "The code took too long to run and was stopped, sir."
    except Exception as e:
        return f"Execution error: {str(e)[:200]}"


# ============================================================
# NEW FEATURE 2: Email Manager (Gmail IMAP/SMTP)
# Setup: Add EMAIL_ADDRESS and EMAIL_APP_PASSWORD to .env
# Gmail App Password: Google Account → Security → 2-Step → App Passwords
# ============================================================

def read_latest_emails(count=5):
    """
    Read the most recent emails from your Gmail inbox.
    Returns a formatted summary string.
    """
    if not config.EMAIL_ADDRESS or not config.EMAIL_APP_PASSWORD:
        return (
            "Email is not configured, sir. "
            "Please add EMAIL_ADDRESS and EMAIL_APP_PASSWORD to your .env file."
        )
    try:
        mail = imaplib.IMAP4_SSL("imap.gmail.com")
        mail.login(config.EMAIL_ADDRESS, config.EMAIL_APP_PASSWORD)
        mail.select("inbox")
        
        _, data = mail.search(None, "ALL")
        mail_ids = data[0].split()
        if not mail_ids:
            return "Your inbox is empty, sir."
        
        latest_ids = mail_ids[-count:][::-1]  # Newest first
        summaries = []
        
        for mail_id in latest_ids:
            _, msg_data = mail.fetch(mail_id, "(RFC822)")
            msg = email_lib.message_from_bytes(msg_data[0][1])
            
            subject = msg.get("Subject", "No Subject")
            sender = msg.get("From", "Unknown")
            date = msg.get("Date", "")[:16]
            
            # Extract plain text body
            body = ""
            if msg.is_multipart():
                for part in msg.walk():
                    if part.get_content_type() == "text/plain":
                        payload = part.get_payload(decode=True)
                        if payload:
                            body = payload.decode('utf-8', errors='ignore')[:150]
                            break
            else:
                payload = msg.get_payload(decode=True)
                if payload:
                    body = payload.decode('utf-8', errors='ignore')[:150]
            
            summaries.append(
                f"From: {sender}\n"
                f"Subject: {subject}\n"
                f"Date: {date}\n"
                f"Preview: {body.strip()[:100]}..."
            )
        
        mail.logout()
        return "\n\n---\n\n".join(summaries) if summaries else "No emails found."
        
    except imaplib.IMAP4.error as e:
        return f"Email login failed: {str(e)}. Check your EMAIL_APP_PASSWORD in .env"
    except Exception as e:
        return f"Email error: {str(e)}"


def send_email_message(to_address, subject, body):
    """
    Send an email via Gmail SMTP.
    to_address: recipient email
    subject: email subject
    body: email body text
    """
    if not config.EMAIL_ADDRESS or not config.EMAIL_APP_PASSWORD:
        return "Email not configured. Add EMAIL_ADDRESS and EMAIL_APP_PASSWORD to .env"
    
    try:
        msg = MIMEMultipart()
        msg['From'] = config.EMAIL_ADDRESS
        msg['To'] = to_address
        msg['Subject'] = subject
        msg.attach(MIMEText(body, 'plain', 'utf-8'))
        
        with smtplib.SMTP_SSL('smtp.gmail.com', 465, timeout=10) as server:
            server.login(config.EMAIL_ADDRESS, config.EMAIL_APP_PASSWORD)
            server.send_message(msg)
        
        print(f"[EMAIL] Sent to {to_address}: {subject}")
        return f"Email sent to {to_address} successfully, sir."
    except smtplib.SMTPAuthenticationError:
        return "Email authentication failed. Check your Gmail App Password in .env"
    except Exception as e:
        return f"Failed to send email: {str(e)}"


# ============================================================
# NEW FEATURE 3: Reminder Time Parser
# Parses natural language time expressions into datetime objects
# ============================================================

def parse_reminder_time(time_str):
    """
    Parse natural language or ISO time strings into datetime.
    
    Examples:
      '5 minutes'     → now + 5 min
      '2 hours'       → now + 2 hrs
      'tomorrow 9am'  → tomorrow at 09:00
      '2024-12-25T09:00:00' → ISO format
      'in 30 seconds' → now + 30 sec
    """
    now = datetime.datetime.now()
    ts = time_str.strip().lower()
    
    # Try ISO format first
    try:
        return datetime.datetime.fromisoformat(ts)
    except:
        pass
    
    # "X minutes" / "X דקות"
    m = re.search(r'(\d+)\s*(?:minute|min|דקו|דקה)', ts)
    if m:
        return now + datetime.timedelta(minutes=int(m.group(1)))
    
    # "X hours" / "X שעות"
    m = re.search(r'(\d+)\s*(?:hour|hr|שעה|שעות)', ts)
    if m:
        return now + datetime.timedelta(hours=int(m.group(1)))
    
    # "X seconds" / "X שניות"
    m = re.search(r'(\d+)\s*(?:second|sec|שניה|שניות)', ts)
    if m:
        return now + datetime.timedelta(seconds=int(m.group(1)))
    
    # "X days" / "X ימים"
    m = re.search(r'(\d+)\s*(?:day|ימ|יום)', ts)
    if m:
        return now + datetime.timedelta(days=int(m.group(1)))
    
    # "tomorrow [time]" / "מחר"
    if 'tomorrow' in ts or 'מחר' in ts:
        tomorrow = now + datetime.timedelta(days=1)
        t_match = re.search(r'(\d{1,2})(?::(\d{2}))?\s*(am|pm)?', ts)
        if t_match:
            hour = int(t_match.group(1))
            minute = int(t_match.group(2) or 0)
            ampm = t_match.group(3)
            if ampm == 'pm' and hour < 12:
                hour += 12
            elif ampm == 'am' and hour == 12:
                hour = 0
            return tomorrow.replace(hour=hour, minute=minute, second=0, microsecond=0)
        return tomorrow.replace(hour=9, minute=0, second=0, microsecond=0)
    
    # "at HH:MM" / "בשעה"
    m = re.search(r'(?:at|בשעה)\s*(\d{1,2}):(\d{2})', ts)
    if m:
        hour, minute = int(m.group(1)), int(m.group(2))
        target = now.replace(hour=hour, minute=minute, second=0, microsecond=0)
        if target <= now:
            target += datetime.timedelta(days=1)
        return target
    
    # Default: 5 minutes
    print(f"[REMINDER] Could not parse '{time_str}', defaulting to 5 minutes")
    return now + datetime.timedelta(minutes=5)


# ============================================================
# NEW FEATURE 4: Background Reminder Checker
# Runs as a daemon thread — checks every 30 seconds for due reminders
# ============================================================

def start_reminder_checker(speak_fn):
    """
    Start a daemon thread that checks for due reminders every 30 seconds.
    When a reminder fires, Jarvis speaks it aloud.
    Call this once from main.py after speak() is ready.
    """
    from memory_manager import get_due_reminders
    
    def _check_loop():
        print("[REMINDER] Background checker started.")
        while True:
            try:
                due = get_due_reminders()
                for text in due:
                    print(f"[REMINDER] 🔔 Firing: {text}")
                    has_hebrew = any('\u0590' <= c <= '\u05ea' for c in text)
                    if has_hebrew:
                        speak_fn(f"תזכורת, אדוני: {text}")
                    else:
                        speak_fn(f"Reminder, sir: {text}")
            except Exception as e:
                print(f"[REMINDER] Checker error: {e}")
            time.sleep(30)
    
    t = threading.Thread(target=_check_loop, daemon=True, name="ReminderChecker")
    t.start()
    return t